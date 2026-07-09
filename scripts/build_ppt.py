"""Build the final UniTok Jittor defense deck from the provided template.

The generated slides follow the advisor's revised requirement:
- concise main view: keywords, diagrams, tables, figures, conclusion lines;
- detailed narration goes to Markdown and a JSON notes file for PowerPoint notes;
- an independent 16-page Jittor code section, 40 seconds per page, about 10 min.
"""

from __future__ import annotations

import csv
import json
import re
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT.parent / "TemplateMC" / "TemplateMC-PPT.pptx"
OUT_DIR = ROOT / "outputs" / "ppt"
PPTX_OUT = OUT_DIR / "UniTok_Jittor_Reproduction.pptx"
SPEECH_OUT = OUT_DIR / "UniTok_Jittor_Speech.md"
NOTES_JSON_OUT = OUT_DIR / "UniTok_Jittor_Speech_Notes.json"
PAPER_PDF = ROOT.parent / "2502.20321v3.pdf"
OFFICIAL_ASSET_DIR = ROOT.parent / "UniTok" / "assets"
PAPER_TEASER = OFFICIAL_ASSET_DIR / "teaser.png"
PAPER_SAMPLES = OFFICIAL_ASSET_DIR / "samples.png"
MODEL_PART_DIR = OUT_DIR / "figures" / "model_parts"

MODEL_PART_CROPS = {
    "overview": (0, 0, 2653, 838),
    "encoder": (420, 120, 910, 560),
    "mcq": (1030, 125, 1660, 545),
    "attention": (760, 115, 1935, 555),
    "decoder": (1620, 125, 2255, 560),
    "loss": (620, 0, 2050, 820),
}

BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(71, 133, 179)
GREEN = RGBColor(84, 130, 53)
ORANGE = RGBColor(198, 89, 17)
DARK = RGBColor(28, 28, 28)
GREY = RGBColor(96, 96, 96)
LIGHT = RGBColor(246, 249, 252)
WHITE = RGBColor(255, 255, 255)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Calibri"

NOTES: list[dict[str, str]] = []


def clear_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for idx in range(len(sld_id_lst) - 1, -1, -1):
        r_id = sld_id_lst[idx].rId
        prs.part.drop_rel(r_id)
        del sld_id_lst[idx]


def style_text_frame(text_frame, size=22, color=DARK, bold=False, align=None) -> None:
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in text_frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = FONT_CN
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size=22, color=DARK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text = text
    style_text_frame(box.text_frame, size, color, bold, align)
    return box


def add_title(slide, title: str) -> None:
    add_text(slide, title, 0.28, 0.16, 12.2, 0.55, size=25, color=DARK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.79), Inches(13.33), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_footer(slide, page_no: int) -> None:
    y = 7.08
    parts = [
        (0.0, 4.0, LIGHT_BLUE),
        (4.0, 8.05, BLUE),
        (12.05, 1.28, RGBColor(10, 45, 82)),
    ]
    for x, w, color in parts:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
    add_text(slide, "jittor-unitok 复现", 0.35, y + 0.08, 3.3, 0.25, size=10, color=WHITE)
    add_text(slide, "UniTok 统一视觉 Tokenizer", 4.65, y + 0.08, 6.7, 0.25, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, str(page_no), 12.42, y + 0.08, 0.45, 0.25, size=10, color=WHITE, align=PP_ALIGN.CENTER)


def new_slide(prs: Presentation, title: str, core: str, body: str):
    page_no = len(prs.slides) + 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)
    add_title(slide, title)
    add_footer(slide, page_no)
    NOTES.append({"page": page_no, "title": title, "core": core, "body": body})
    return slide


def add_bullets(slide, bullets: Iterable[str], x=0.85, y=1.35, w=11.7, h=4.7, size=22) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = FONT_CN
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_tag(slide, text: str, x: float, y: float, w: float, color=BLUE, size=13) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    add_text(slide, text, x + 0.08, y + 0.08, w - 0.16, 0.23, size=size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, color=BLUE, body_size=18) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT
    rect.line.color.rgb = color
    rect.line.width = Pt(1.2)
    add_text(slide, title, x + 0.16, y + 0.14, w - 0.32, 0.34, size=16, color=color, bold=True)
    add_text(slide, body, x + 0.16, y + 0.58, w - 0.32, h - 0.7, size=body_size, color=DARK)


def add_conclusion(slide, text: str, y=6.15) -> None:
    add_tag(slide, text, 1.0, y, 11.3, ORANGE, size=14)


def add_code(slide, code: str, x: float, y: float, w: float, h: float, size=11) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(247, 247, 247)
    rect.line.color.rgb = RGBColor(180, 180, 180)
    box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.1), Inches(w - 0.24), Inches(h - 0.18))
    box.text = code
    box.text_frame.word_wrap = False
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    for paragraph in box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Consolas"
            run.font.size = Pt(size)
            run.font.color.rgb = DARK


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, font_size=14) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else RGBColor(250, 250, 250)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT_CN
                    run.font.size = Pt(font_size)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else DARK


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color=BLUE) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True


def add_flow(slide, labels: list[str], y=3.0, x0=0.65, box_w=1.55, gap=0.36, color=BLUE) -> None:
    for i, label in enumerate(labels):
        x = x0 + i * (box_w + gap)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.72))
        rect.fill.solid()
        rect.fill.fore_color.rgb = LIGHT
        rect.line.color.rgb = color
        add_text(slide, label, x + 0.05, y + 0.18, box_w - 0.1, 0.28, size=13, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_arrow(slide, x + box_w, y + 0.36, x + box_w + gap - 0.05, y + 0.36, color)


def picture(slide, path: Path, x: float, y: float, w: float) -> bool:
    if not path.exists():
        return False
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return True


def ensure_model_part_assets() -> None:
    """Crop official UniTok teaser into small model-location images for code slides."""

    if not PAPER_TEASER.exists():
        return
    MODEL_PART_DIR.mkdir(parents=True, exist_ok=True)
    image = Image.open(PAPER_TEASER).convert("RGB")
    for name, box in MODEL_PART_CROPS.items():
        out = MODEL_PART_DIR / f"{name}.png"
        image.crop(box).save(out)


def picture_fit(slide, path: Path, x: float, y: float, w: float, h: float) -> bool:
    """Insert image into a fixed box while preserving aspect ratio."""

    if not path.exists():
        return False
    with Image.open(path) as image:
        img_w, img_h = image.size
    scale = min(w / img_w, h / img_h)
    draw_w = img_w * scale
    draw_h = img_h * scale
    px = x + (w - draw_w) / 2
    py = y + (h - draw_h) / 2
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(250, 250, 250)
    frame.line.color.rgb = RGBColor(190, 190, 190)
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(draw_w), height=Inches(draw_h))
    return True


def paper_picture(slide, path: Path, x: float, y: float, w: float, label: str) -> None:
    """Insert an official paper asset, or reserve a visible slot if it is missing."""

    if picture(slide, path, x, y, w):
        add_text(slide, f"图源：{label}", x, y + 3.95, w, 0.25, size=10, color=GREY)
        return
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(3.6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(250, 250, 250)
    rect.line.color.rgb = ORANGE
    add_text(slide, "原论文模型结构图预留位", x + 0.3, y + 1.35, w - 0.6, 0.45, size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"请放入：{label}", x + 0.3, y + 2.0, w - 0.6, 0.35, size=16, color=GREY, align=PP_ALIGN.CENTER)


def read_loss_tail() -> str:
    csv_path = ROOT / "outputs" / "curves" / "loss.csv"
    if not csv_path.exists():
        return "loss.csv 未生成"
    rows = list(csv.DictReader(csv_path.open("r", encoding="utf-8")))
    if not rows:
        return "loss.csv 为空"
    row = rows[-1]
    return f"step {row['step']} | total {float(row['total_loss']):.3f} | recon {float(row['recon_loss']):.3f} | vq {float(row['vq_loss']):.3f}"


def read_epoch_loss_summary() -> str:
    csv_path = ROOT / "outputs" / "curves" / "loss.csv"
    if not csv_path.exists():
        return "loss.csv 未生成"
    rows = list(csv.DictReader(csv_path.open("r", encoding="utf-8")))
    buckets: dict[int, list[float]] = {}
    for row in rows:
        epoch = int(float(row["epoch"]))
        buckets.setdefault(epoch, []).append(float(row["total_loss"]))
    first_epoch = min(buckets)
    last_epoch = max(buckets)
    first_mean = sum(buckets[first_epoch]) / len(buckets[first_epoch])
    last_mean = sum(buckets[last_epoch]) / len(buckets[last_epoch])
    drop = (first_mean - last_mean) / max(abs(first_mean), 1e-8) * 100.0
    return f"epoch {first_epoch}: {first_mean:.3f}\nepoch {last_epoch}: {last_mean:.3f}\n下降 {drop:.1f}%"


def read_log_summary() -> str:
    log_path = ROOT / "outputs" / "logs" / "train.log"
    if not log_path.exists():
        return "train.log 未生成"
    lines = [x.strip() for x in log_path.read_text(encoding="utf-8", errors="ignore").splitlines() if "epoch=" in x or "saved" in x]
    return "\n".join(lines[-3:])


def read_final_train_summary() -> str:
    log_path = ROOT / "outputs" / "logs" / "train.log"
    if not log_path.exists():
        return "train.log 未生成"
    pattern = re.compile(
        r"epoch=(?P<epoch>\d+)\s+step=(?P<step>\d+)\s+total_loss=(?P<total>[0-9.]+)\s+"
        r"recon_loss=(?P<recon>[0-9.]+)\s+vq_loss=(?P<vq>[0-9.]+)"
    )
    for line in reversed(log_path.read_text(encoding="utf-8", errors="ignore").splitlines()):
        match = pattern.search(line)
        if match:
            return (
                f"train.log: epoch {match.group('epoch')} | step {match.group('step')}\n"
                f"total {float(match.group('total')):.3f} | recon {float(match.group('recon')):.3f} | vq {float(match.group('vq')):.3f}"
            )
    return "train.log 未找到 epoch 行"


def read_metrics() -> dict[str, float]:
    path = ROOT / "outputs" / "logs" / "eval_metrics.json"
    if not path.exists():
        return {"mse": 0.0, "l1": 0.0, "psnr": 0.0}
    return json.loads(path.read_text(encoding="utf-8"))


def read_training_manifest() -> str:
    path = ROOT / "outputs" / "ppt" / "figures" / "manifest.json"
    if not path.exists():
        return "manifest.json 未生成"
    manifest = json.loads(path.read_text(encoding="utf-8"))
    return f"epoch {manifest.get('final_epoch', '?')} | step {manifest.get('final_step', '?')} | MA {manifest.get('moving_average_window', '?')}"


CODE_SLIDES = [
    {
        "title": "代码 1/16：项目目录结构",
        "path": "jittor-unitok/",
        "module": "论文对应：tokenizer 训练工程",
        "shape": "输入输出：项目级组织，无张量输入",
        "solves": "把论文算法拆成可运行、可训练、可展示的工程闭环。",
        "code": """jittor_unitok/
  models/    # 核心模型：MCQ、注意力投影、Tokenizer、Loss
  engine/    # 训练与评估入口：train / eval
scripts/     # 数据准备、曲线绘制、PyTorch/Jittor 对齐
tests/       # shape、反向传播、单步训练测试
outputs/     # 日志、曲线、重建图、PPT 成品""",
        "notes": "代码讲解第 1 页，预计用时 40 秒。这里先看项目整体结构。models 目录承载 UniTok tokenizer 的核心算法，engine 负责训练和评估入口，scripts 负责数据、曲线、PyTorch/Jittor 对齐和 PPT 生成，outputs 保存所有可提交产物。这个结构的意义是让论文复现从模型实现走到实验证据，而不是只写孤立模块。",
    },
    {
        "title": "代码 2/16：tokenizer.py 总流程",
        "path": "jittor_unitok/models/tokenizer.py",
        "module": "论文对应：UniTok tokenizer graph",
        "shape": "image [B,3,H,W] -> indices [B,M,N] -> recon [B,3,H,W]",
        "solves": "串起 encoder、projection、MCQ、decoder 的完整路径。",
        "code": """def execute(self, images):
    # 图像 -> 连续视觉 token: [B,3,H,W] -> [B,N,C]
    tokens = self.encoder(images)
    # 注意力投影压缩通道，降低量化难度
    latent = self.quant_proj(tokens)
    # MCQ 输出量化特征、离散索引和 VQ 训练项
    quantized, indices, commitment, codebook, stats = self.quantizer(latent)
    # 量化特征恢复到 decoder 需要的 hidden_dim
    decoded_tokens = self.post_quant_proj(quantized)
    recon = self.decoder(decoded_tokens)       # [B,3,H,W]
    rec_loss = reconstruction_loss(recon, images, self.recon_loss_type)
    quant_loss = vq_loss(commitment, codebook)""",
        "notes": "代码讲解第 2 页，预计用时 40 秒。本页不再只放流程示意，而是展示 tokenizer.py 的真实 forward 主干。它完整串起 encoder、compression、MCQ、expansion、decoder 和两类损失。重点看 quantizer 返回 commitment、codebook 和 stats，这些正好对应论文中的 VQ 训练和码本统计；最后 pooled visual feature 则保留理解侧接口。",
    },
    {
        "title": "代码 3/16：tokenizer 接口",
        "path": "jittor_unitok/models/tokenizer.py",
        "module": "论文对应：离散视觉 token 接口",
        "shape": "encode: [B,3,H,W] -> [B,M,N]；decode: [B,M,N] -> [B,3,H,W]",
        "solves": "把训练、推理和 MLLM 接入接口分开。",
        "code": """def encode(self, images):
    # 推理时只导出离散 token，不走 decoder
    tokens = self.encoder(images)
    latent = self.quant_proj(tokens)
    return self.quantizer.f_to_idx(latent)     # [B,M,N]

def decode(self, indices):
    # M 个子码本索引查表，恢复量化 feature
    quantized = self.quantizer.idx_to_f(indices)
    decoded_tokens = self.post_quant_proj(quantized)
    return self.decoder(decoded_tokens)        # [B,3,H,W]

def reconstruct(self, images):
    # 评估脚本直接复用完整 forward
    return self.execute(images)[\"recon\"]""",
        "notes": "代码讲解第 3 页，预计用时 40 秒。这一页展示三个真实接口：encode 明确经过 encoder 和 quant_proj，再输出离散 indices；decode 从 indices 查表恢复 quantized feature，再 expansion 和 decoder；reconstruct 则复用 execute。这样老师能看到项目不是只有训练 forward，而是具备真正的离散 token 导出和恢复接口。",
    },
    {
        "title": "代码 4/16：MCQ 输入输出 shape",
        "path": "jittor_unitok/models/mcq.py",
        "module": "论文对应：Multi-codebook Quantization",
        "shape": "features [B,N,C] -> quantized [B,N,C] + indices [B,M,N]",
        "solves": "用多个子码本共同表示一个视觉 token。",
        "code": """def execute(self, features):
    # features: [B,N,C]，C 会被切成 M 个 chunk
    quantized_chunks, index_chunks = [], []
    commitment_loss = jt.zeros((1,), dtype=\"float32\")
    codebook_loss = jt.zeros((1,), dtype=\"float32\")
    for i in range(self.num_codebooks):
        start = i * self.chunk_dim
        end = (i + 1) * self.chunk_dim
        chunk = features[:, :, start:end]      # [B,N,C/M]
        # 第 i 个 chunk 只查第 i 个子码本
        q, idx = self._nearest_code(chunk, self.codebooks[i])
        quantized_chunks.append(q)
        index_chunks.append(idx)""",
        "notes": "代码讲解第 4 页，预计用时 40 秒。本页展示 MCQ execute 的真实开头。可以看到输入 features 不再被一个码本整体量化，而是在循环里根据 num_codebooks 切成多个 chunk。每个 chunk 进入自己的 codebook，并同时收集量化结果和 index。这是 Multi-codebook Quantization 和普通 VQ 的核心差别。",
    },
    {
        "title": "代码 5/16：chunk / lookup / concat",
        "path": "jittor_unitok/models/mcq.py",
        "module": "论文公式：f_i -> Z_i -> z_i",
        "shape": "chunk [B,N,C/M]；concat 后恢复 [B,N,C]",
        "solves": "把高维量化拆成多个低维最近邻问题。",
        "code": """for i in range(self.num_codebooks):
    # 公式 f = concat(f_1, ..., f_M)
    start = i * self.chunk_dim
    end = (i + 1) * self.chunk_dim
    chunk = features[:, :, start:end]      # [B, N, C/M]
    # z_i = nearest_code(f_i, Z_i)
    q, idx = self._nearest_code(chunk, self.codebooks[i])
    quantized_chunks.append(q)
    index_chunks.append(idx)

# hat_f = concat(z_1, ..., z_M)
quantized = jt.concat(quantized_chunks, dim=-1)  # [B, N, C]
indices = jt.stack(index_chunks, dim=1)          # [B, M, N]""",
        "notes": "代码讲解第 5 页，预计用时 40 秒。这里完整展示 split、lookup、concat 的代码段。start/end 决定第 i 个 chunk 的通道范围，chunk 的 shape 是 B、N、C/M。所有子码本输出 q 后，用 concat 恢复 B、N、C；所有 idx 用 stack 得到 B、M、N。这正好对应论文公式 f_i 分别量化后再拼接成 hat_f。",
    },
    {
        "title": "代码 6/16：nearest code lookup",
        "path": "jittor_unitok/models/mcq.py",
        "module": "论文对应：最近邻码字搜索",
        "shape": "chunk [B,N,D] + codebook [K,D] -> idx [B,N]",
        "solves": "把连续子向量映射成离散 code id。",
        "code": """bsz, num_tokens, dim = chunk.shape
flat = chunk.reshape((-1, dim))
# 搜索时归一化，等价于余弦最近邻
flat_for_search = _l2_normalize(flat, dim=-1).stop_grad()
book_for_search = _l2_normalize(codebook, dim=-1).stop_grad()
x2 = (flat_for_search * flat_for_search).sum(dim=1, keepdims=True)
z2 = (book_for_search * book_for_search).sum(dim=1).reshape((1, -1))
dist = x2 + z2 - 2.0 * (flat_for_search @ book_for_search.transpose())
indices, _ = jt.argmin(dist, dim=1)       # 离散 code id
eye = np.eye(self.codebook_size, dtype=\"float32\")
one_hot = jt.array(eye[indices.numpy()])
quantized = one_hot @ codebook""",
        "notes": "代码讲解第 6 页，预计用时 40 秒。这里展示最近邻查找的完整关键路径：先把 B、N、D 的 chunk 展平，再对 feature 和 codebook 做归一化，然后用 x2+z2-2xz 计算 L2 距离。argmin 产生离散 index，one_hot 乘 codebook 完成真正的查表。这页能说明离散 token 是怎样从连续向量来的。",
    },
    {
        "title": "代码 7/16：straight-through estimator",
        "path": "jittor_unitok/models/mcq.py",
        "module": "论文对应：离散量化的反向传播",
        "shape": "features/quantized [B,N,C] -> quantized_st [B,N,C]",
        "solves": "前向使用离散码字，反向仍训练 encoder。",
        "code": """quantized = jt.concat(quantized_chunks, dim=-1)
indices = jt.stack(index_chunks, dim=1)
commitment_loss = commitment_loss / self.num_codebooks
codebook_loss = codebook_loss / self.num_codebooks

# 直通估计：前向用离散码字，反向梯度走 features
quantized_st = features + (quantized - features).stop_grad()
usage = self._usage_stats(indices)
return quantized_st, indices, commitment_loss, codebook_loss, usage""",
        "notes": "代码讲解第 7 页，预计用时 40 秒。这里展示 ST estimator 所在的真实上下文。前面 concat 得到离散量化后的 quantized，indices 作为离散 token 输出；quantized_st 用 features 加 stop_grad 差值实现直通估计。讲这页时要强调：前向使用码字，反向梯度仍能回到 encoder 和 projection。",
    },
    {
        "title": "代码 8/16：VQ loss 与 usage",
        "path": "jittor_unitok/models/mcq.py",
        "module": "论文对应：L_VQ 与码本利用率",
        "shape": "loss 标量；usage/perplexity 标量统计",
        "solves": "约束 latent 靠近码字，并监控码本是否被使用。",
        "code": """# commitment: 约束 encoder latent 靠近选中码字
commitment_loss = commitment_loss + ((q.stop_grad() - chunk) ** 2).mean() * self.beta
# soft codebook loss: 轻量实现中稳定更新码本
codebook_loss = codebook_loss + self._soft_codebook_loss(chunk, self.codebooks[i])

def _usage_stats(self, indices):
    # 统计有多少 code 真正被 batch 使用
    flat = indices.reshape((-1,))
    one_hot = nn.one_hot(flat, self.codebook_size).float32()
    probs = one_hot.mean(dim=0)
    used = (probs > 0).float32().mean()
    entropy = -(probs * jt.log(probs + 1e-7)).sum()
    return {\"usage\": used, \"perplexity\": jt.exp(entropy)}""",
        "notes": "代码讲解第 8 页，预计用时 40 秒。这页展示 VQ loss 和 usage 的真实实现。commitment loss 让 encoder 输出不要离选中码字太远，soft codebook loss 给码本一个稳定更新路径；usage_stats 用 one-hot 统计实际被用到的码字比例和 perplexity。你训练后 usage 接近 1，说明码本没有塌缩。",
    },
    {
        "title": "代码 9/16：ChannelCompressionBlock",
        "path": "jittor_unitok/models/attention_projection.py",
        "module": "论文对应：Attention Projection compression",
        "shape": "[B,N,C] -> [B,N,c]",
        "solves": "在量化前降低通道维度，缓解量化难度。",
        "code": """def execute(self, x):
    # x: [B,N,C]，在 token 维度做自注意力
    bsz, num_tokens, _ = x.shape
    y = self.norm1(x)
    # q/k/v 投影后拆成多头
    q = self.q(y).reshape((bsz, num_tokens, self.num_heads, self.head_dim))
    k = self.k(y).reshape((bsz, num_tokens, self.num_heads, self.head_dim))
    v = self.v(y).reshape((bsz, num_tokens, self.num_heads, self.head_dim))
    q = q.transpose((0, 2, 1, 3))
    k = k.transpose((0, 2, 1, 3))
    v = v.transpose((0, 2, 1, 3))
    # attention 保留 token 间上下文
    scores = (q @ k.transpose((0,1,3,2))) * self.scale
    attn = nn.softmax(scores, dim=-1)
    y = attn @ v
    y = y.transpose((0, 2, 1, 3))
    return y.reshape((bsz, num_tokens, self.out_dim))""",
        "notes": "代码讲解第 9 页，预计用时 40 秒。这页展示 Attention Projection 的核心 execute，而不只是类名。q、k、v 都从输入 token 投影出来，再 reshape 成多头形式，在 token 维度上计算 attention。CompressionBlock 继承这套逻辑，只是把 out_dim 设置成 compressed_dim，因此实现的是 N、C 到 N、c 的注意力压缩。",
    },
    {
        "title": "代码 10/16：ChannelExpansionBlock",
        "path": "jittor_unitok/models/attention_projection.py",
        "module": "论文对应：Attention Projection expansion",
        "shape": "[B,N,c] -> [B,N,C]",
        "solves": "在解码前恢复 decoder 需要的 hidden_dim。",
        "code": """class ChannelCompressionBlock(ChannelAttentionProjection):
    # 量化前：hidden_dim -> latent_dim
    def __init__(self, in_dim, compressed_dim, num_heads=4):
        super().__init__(in_dim=in_dim,
                         out_dim=compressed_dim,
                         num_heads=num_heads)

class ChannelExpansionBlock(ChannelAttentionProjection):
    # 解码前：latent_dim -> hidden_dim
    def __init__(self, compressed_dim, out_dim, num_heads=4):
        super().__init__(in_dim=compressed_dim,
                         out_dim=out_dim,
                         num_heads=num_heads)""",
        "notes": "代码讲解第 10 页，预计用时 40 秒。这页展示 compression 和 expansion 是同一个 ChannelAttentionProjection 的两个实例化方向。Compression 把 hidden_dim 压到 latent_dim，Expansion 把 latent_dim 恢复到 hidden_dim。这样实现上很紧凑，但论文语义上对应量化前后的通道 factorization 闭环。",
    },
    {
        "title": "代码 11/16：tiny backbone",
        "path": "jittor_unitok/models/encoder.py, decoder.py",
        "module": "论文对应：ViTamin-L/16 的轻量替代",
        "shape": "encoder [B,3,H,W] -> [B,N,C]；decoder [B,N,C] -> [B,3,H,W]",
        "solves": "在普通机器上验证 tokenizer 链路。",
        "code": """# encoder.py
feat = self.net(x)                         # [B, C, H/8, W/8]
bsz, channels, height, width = feat.shape
# 2D feature map 展平成 token 序列
tokens = feat.reshape((bsz, channels, height * width))
tokens = tokens.transpose((0, 2, 1))

# decoder.py
bsz, num_tokens, channels = tokens.shape
expected = self.grid_size * self.grid_size
# token 序列恢复成 2D feature map
feat = tokens.transpose((0, 2, 1))
feat = feat.reshape((bsz, channels, self.grid_size, self.grid_size))
recon = jt.tanh(self.net(feat))            # [B, 3, H, W]""",
        "notes": "代码讲解第 11 页，预计用时 40 秒。这里展示 encoder 和 decoder 的真实 shape 变换。encoder 先输出二维 feature map，再 flatten 成 token 序列；decoder 则检查 token 数量，把序列 reshape 回 grid，再上采样重建图像。原论文 backbone 更大，但这个 tiny backbone 保留了 image-token-image 的核心接口。",
    },
    {
        "title": "代码 12/16：重建损失与 VQ 损失",
        "path": "jittor_unitok/models/losses.py",
        "module": "论文对应：L_R + L_VQ",
        "shape": "pred/target [B,3,H,W]；commitment/codebook 为标量",
        "solves": "把论文完整目标中的核心重建项落实为 tiny 可训练损失。",
        "font_size": 11.0,
        "code": """def reconstruction_loss(pred, target, loss_type=\"l1\"):
    # L_R: decoder 输出与原图之间的像素误差
    if loss_type == \"l1\":
        return jt.abs(pred - target).mean()
    # MSE 用于和评估指标 MSE/PSNR 对齐
    if loss_type == \"mse\":
        return ((pred - target) ** 2).mean()
    raise ValueError(f\"Unsupported loss: {loss_type}\")

def vq_loss(commitment_loss, codebook_loss):
    # L_VQ = commitment loss + codebook loss
    return commitment_loss + codebook_loss""",
        "notes": "代码讲解第 12 页，预计用时 40 秒。这一页完整展示 losses.py 中重建损失和 VQ 损失的定义。reconstruction_loss 对应论文里的 L_R，输入是 decoder 输出 pred 和原图 target，shape 都是 B、3、H、W；默认 L1 更稳定，MSE 分支用于和评估指标对齐。vq_loss 对应 L_VQ，本项目把 commitment loss 和 codebook loss 相加，lambda_VQ 在训练循环外统一加权。",
    },
    {
        "title": "代码 13/16：对比损失接口",
        "path": "jittor_unitok/models/losses.py",
        "module": "论文对应：L_contra",
        "shape": "image/text feature [B,D] -> 相似度矩阵 [B,B] -> loss 标量",
        "solves": "保留理解侧图文对齐接口，但 tiny 默认不强制启用。",
        "font_size": 10.6,
        "code": """def contrastive_loss(image_features, text_features, temperature=0.07):
    # L2 normalize 后，点积就是余弦相似度
    image_norm = jt.sqrt((image_features * image_features)
                         .sum(dim=-1, keepdims=True) + 1e-6)
    text_norm = jt.sqrt((text_features * text_features)
                        .sum(dim=-1, keepdims=True) + 1e-6)
    image_features = image_features / image_norm
    text_features = text_features / text_norm

    logits = image_features @ text_features.transpose() / temperature
    labels = jt.arange(logits.shape[0])  # 对角线是正样本
    loss_i2t = nn.cross_entropy_loss(logits, labels)
    loss_t2i = nn.cross_entropy_loss(logits.transpose(), labels)
    return 0.5 * (loss_i2t + loss_t2i)""",
        "notes": "代码讲解第 13 页，预计用时 40 秒。这一页完整展示 contrastive_loss。它先把图像特征和文本特征做 L2 normalize，然后构造 B 乘 B 的相似度矩阵 logits，矩阵对角线是正样本。image-to-text 和 text-to-image 两个方向分别做交叉熵，最后取平均。这里和论文 L_contra 对应，但本项目没有大规模真实 caption，所以默认训练不启用，只保留接口和可扩展性。",
    },
    {
        "title": "代码 14/16：训练循环与产物",
        "path": "jittor_unitok/engine/train_tokenizer.py",
        "module": "论文对应：tokenizer training",
        "shape": "batch image [B,3,H,W] -> loss 标量 + checkpoint/PNG/CSV",
        "solves": "把算法训练变成可复现的实验记录。",
        "code": """for images_np, names in iter_image_batches(train_set, cfg[\"batch_size\"], shuffle=True):
    images = jt.array(images_np)            # [B,3,H,W]
    out = model(images)
    # tiny 默认目标：L_R + lambda_vq * L_VQ
    total_loss = out[\"recon_loss\"] + cfg[\"lambda_vq\"] * out[\"vq_loss\"]
    optimizer.step(total_loss)
    global_step += 1
    # 每 step 写 CSV，供 loss 曲线脚本读取
    row = {\"epoch\": epoch, \"step\": global_step,
           \"total_loss\": scalar(total_loss),
           \"recon_loss\": scalar(out[\"recon_loss\"]),
           \"vq_loss\": scalar(out[\"vq_loss\"])}
    append_csv(curve_dir / \"loss.csv\", row, CSV_FIELDS)""",
        "notes": "代码讲解第 14 页，预计用时 40 秒。这页展示真实训练循环：batch 进入模型，total_loss 由重建损失和 VQ 损失组成，然后 optimizer.step 更新参数。row 会写入 loss.csv，这是后面所有 loss 曲线的来源。这样 PPT 中的训练图不是手工造的，而是由这个循环逐 step 记录出来的。",
    },
    {
        "title": "代码 15/16：评估与 loss 曲线",
        "path": "eval_reconstruction.py, scripts/plot_loss.py",
        "module": "论文对应：重建评估与训练趋势",
        "shape": "recon/image [B,3,H,W] -> MSE/L1/PSNR；CSV -> PNG",
        "solves": "用轻量指标证明 encode-decode 和训练趋势有效。",
        "code": """# eval_reconstruction.py
# 测试 encode -> decode 是否能重建图像
recon = model.reconstruct(images)
metrics = reconstruction_metrics(recon, images)
metric_sum[\"mse\"] += metrics[\"mse\"]
save_reconstruction_grid(images, recon, recon_dir / f\"eval_batch_{batches}.png\", names)

# plot_loss.py
# 读取训练 CSV，绘制 total/recon/VQ 三条曲线
rows = list(csv.DictReader(open(args.csv, \"r\", encoding=\"utf-8\")))
steps = [int(r[\"step\"]) for r in rows]
for key in (\"total_loss\", \"recon_loss\", \"vq_loss\"):
    plt.plot(steps, [float(r[key]) for r in rows], label=key)""",
        "notes": "代码讲解第 15 页，预计用时 40 秒。这里把评估和曲线生成都放出关键代码。eval_reconstruction 用 model.reconstruct 得到重建图，再计算 MSE、L1、PSNR 并保存网格图；plot_loss.py 读取同一个 loss.csv，把 total、recon、VQ 三条曲线画出来。它们共同构成结果展示页的数据来源。",
    },
    {
        "title": "代码 16/16：对齐与一键复现",
        "path": "scripts/compare_with_pytorch.py, README.md",
        "module": "论文对应：PyTorch/Jittor 复现对齐",
        "shape": "模块表 + shape 表 + loss 趋势记录",
        "solves": "说明 Jittor 版本对齐的是官方结构和可运行行为。",
        "code": """report = {
  # 对齐官方 PyTorch 模块与 Jittor 文件
  \"module_alignment\": [
    [\"models/unitok.py\", \"jittor_unitok/models/tokenizer.py\", \"tokenizer graph\"],
    [\"models/quant.py\", \"jittor_unitok/models/mcq.py\", \"MCQ\"],
    [\"models/vqvae.py\", \"jittor_unitok/models/attention_projection.py\", \"attention\"],
  ],
  # 对齐关键输入输出 shape
  \"shape_alignment\": {
    \"image\": \"[B,3,H,W]\",
    \"indices\": \"[B,num_codebooks,N]\",
    \"reconstruction\": \"[B,3,H,W]\",
  },
}""",
        "notes": "代码讲解第 16 页，预计用时 40 秒。这页展示 compare_with_pytorch.py 真实输出结构：一部分是模块职责对齐，另一部分是关键 shape 对齐。它说明 Jittor 版本不是随便写了一个 VAE，而是按官方 UniTok 的 tokenizer、quantizer、attention projection 三条主线做了结构级迁移。",
    },
]


def add_code_slide(prs: Presentation, spec: dict) -> None:
    core = f"{spec['title']}：{spec['solves']}"
    slide = new_slide(prs, spec["title"], core, spec["notes"])
    visual_key = code_visual_key(spec)
    add_tag(slide, f"文件路径：{spec['path']}", 0.52, 1.02, 4.75, BLUE, size=11)
    add_tag(slide, spec["module"], 5.42, 1.02, 4.1, GREEN, size=11)
    add_tag(slide, "对应模型位置", 9.82, 1.02, 2.55, ORANGE, size=11)
    add_code(slide, spec["code"], 0.48, 1.48, 9.0, 5.12, size=spec.get("font_size", 11.6))
    model_part = MODEL_PART_DIR / f"{visual_key}.png"
    if not picture_fit(slide, model_part, 9.78, 1.52, 2.7, 1.4):
        add_card(slide, "模型位置", spec["module"], 9.78, 1.52, 2.7, 1.4, ORANGE, body_size=11)
    add_card(slide, "输入 / 输出", spec["shape"], 9.78, 3.12, 2.7, 1.35, BLUE, body_size=11)
    add_card(slide, "作用", spec["solves"], 9.78, 4.72, 2.7, 1.42, ORANGE, body_size=11)


def code_visual_key(spec: dict) -> str:
    path = spec["path"]
    title = spec["title"]
    if "mcq.py" in path:
        return "mcq"
    if "attention_projection.py" in path:
        return "attention"
    if "encoder.py" in path or "decoder.py" in path:
        return "encoder" if "tiny backbone" in title else "decoder"
    if "losses.py" in path:
        return "loss"
    if "train_tokenizer.py" in path or "eval_reconstruction.py" in path or "plot_loss.py" in path:
        return "loss"
    return "overview"


def build_deck() -> Presentation:
    NOTES.clear()
    ensure_model_part_assets()
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(3.75))
    top.fill.solid()
    top.fill.fore_color.rgb = BLUE
    top.line.fill.background()
    add_text(slide, "UniTok Jittor 复现", 1.0, 0.95, 11.3, 0.7, size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "统一视觉生成与理解 Tokenizer 的核心结构复现", 1.5, 1.9, 10.3, 0.45, size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "多码本量化 · 注意力投影 · 小规模训练 · 可视化", 1.9, 2.65, 9.6, 0.35, size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "基于项目模板制作 | jittor-unitok", 3.0, 4.75, 7.3, 0.45, size=22, color=DARK, align=PP_ALIGN.CENTER)
    add_footer(slide, 1)
    NOTES.append({"page": 1, "title": "标题", "core": "本项目复现 UniTok tokenizer 的核心结构，而不是完整大规模训练。", "body": "大家好，这次汇报的主题是 UniTok 的 Jittor 复现。我的复现重点不是重新训练论文里的大模型，而是把统一视觉 tokenizer 的核心结构迁移到 Jittor，包括 MCQ 多码本量化、Attention Projection、forward/encode/decode、训练日志、loss 曲线和重建可视化。"})

    slide = new_slide(prs, "故事线", "从统一 tokenizer 的必要性讲到 Jittor tiny 复现证据。", "这一页给出整份汇报的故事线。先从问题出发：为什么视觉生成和理解需要统一 tokenizer。然后看现有方法：CLIP 擅长理解，VQVAE/VQGAN 擅长生成，但两者直接相加不是最优。核心矛盾是离散 token 表达能力不足。接下来讲 UniTok 的 MCQ 和 Attention Projection，再进入 Jittor 代码详解，最后展示 tiny setting 下的训练、重建和对齐结果。")
    add_flow(slide, ["问题", "现有方法", "核心矛盾", "UniTok 方法", "Jittor 复现", "结果", "思考"], y=2.65, x0=0.85, box_w=1.45, gap=0.35)
    add_conclusion(slide, "主线：离散 token 容量决定统一生成与理解的上限")

    slide = new_slide(prs, "问题：为什么需要统一 tokenizer？", "统一视觉 token 是生成、理解和 MLLM 接口之间的共同语言。", "视觉生成需要能被 decoder 使用的离散 token，视觉理解需要能和文本语义对齐的表示。如果这两套 token 完全分开，系统会复杂，也很难接入统一多模态模型。统一 tokenizer 的目标是让一套视觉 token 同时支持生成和理解。")
    add_card(slide, "生成侧", "需要可解码\n保留视觉细节", 1.0, 1.65, 3.2, 2.1, BLUE)
    add_card(slide, "理解侧", "需要语义对齐\n支持图文任务", 5.05, 1.65, 3.2, 2.1, GREEN)
    add_card(slide, "统一接口", "离散视觉 token\n接入 MLLM", 9.1, 1.65, 3.2, 2.1, ORANGE)
    add_conclusion(slide, "核心问题：一个 token 空间能否同时服务生成和理解？")

    slide = new_slide(prs, "现有方法：各有优势，也各有短板", "CLIP 强在语义，VQVAE/VQGAN 强在生成，但能力没有自然统一。", "CLIP 的优势是图文语义对齐，适合理解任务；但它不是为离散图像生成设计的。VQVAE 和 VQGAN 的优势是离散重建和生成，但语义理解通常比较弱。UniTok 的动机就是把这两类能力放进同一个 tokenizer。")
    add_table(slide, [["方法", "优势", "短板"], ["CLIP", "图文语义强", "不适合离散生成"], ["VQVAE/VQGAN", "重建/生成友好", "语义理解弱"], ["简单拼接", "看似兼顾", "目标冲突 + 容量瓶颈"]], 1.0, 1.45, 11.3, 3.6, font_size=16)
    add_conclusion(slide, "简单 VQVAE + CLIP 不是最优解")

    slide = new_slide(prs, "核心矛盾：Quantization Bottleneck", "真正瓶颈不是 loss 数量，而是离散 token 表达能力。", "如果连续视觉特征很丰富，但离散码本容量不足，那么无论加多少监督，token 都可能表达不完整。码本利用率低、语义和细节争夺容量，都会限制统一 tokenizer 的效果。")
    add_flow(slide, ["连续视觉特征", "有限 codebook", "离散 token", "生成 + 理解"], y=2.25, x0=1.4, box_w=2.0, gap=0.55)
    add_card(slide, "瓶颈表现", "码本利用率低\n细节与语义争夺容量\n高维最近邻困难", 2.0, 4.1, 4.0, 1.45, ORANGE)
    add_card(slide, "关键判断", "需要提升 token 组合表达能力", 7.1, 4.1, 4.0, 1.45, GREEN)
    add_conclusion(slide, "UniTok 从 tokenizer 容量角度解决统一表示问题")

    slide = new_slide(prs, "方法：MCQ + Attention Projection", "UniTok 用多码本和注意力投影增强离散 token 表达能力。", "UniTok 的核心不是简单叠加 VQVAE loss 和 CLIP loss，而是改变 tokenizer 结构。MCQ 把一个 token 拆成多个子空间量化，扩大组合容量；Attention Projection 在压缩和恢复通道时保留 token 间上下文。")
    add_card(slide, "MCQ", "多子码本\n组合式离散表达", 1.2, 1.6, 4.6, 2.1, BLUE)
    add_card(slide, "Attention Projection", "N×C ⇄ N×c\n注意力压缩/恢复", 7.5, 1.6, 4.6, 2.1, GREEN)
    add_flow(slide, ["token", "split", "lookup", "concat", "decode"], y=4.55, x0=2.0, box_w=1.55, gap=0.45, color=ORANGE)
    add_conclusion(slide, "结构改造优先于单纯堆叠训练目标")

    slide = new_slide(prs, "原论文模型图：统一 token 的位置", "UniTok 的关键结构位于视觉编码器和解码器之间。", "这一页引用原论文和官方仓库中的模型图。按图从左到右看，图像先进入视觉编码器，得到连续视觉 token；中间的 Attention Projection 先做通道压缩，随后 MCQ 多个 VQ 分支分别量化子空间；右侧 Attention Projection 再把量化特征恢复到解码器需要的维度。下方文本编码器和对比损失说明同一套视觉 token 还要服务理解任务。")
    paper_picture(slide, PAPER_TEASER, 0.55, 1.05, 12.2, "原论文 PDF / UniTok/assets/teaser.png")
    add_tag(slide, "主线：图像编码 -> 注意力压缩 -> MCQ -> 注意力恢复 -> 图像重建", 1.0, 5.32, 11.3, BLUE, size=13)
    add_conclusion(slide, "MCQ 和 Attention Projection 是统一 tokenizer 的结构核心")

    slide = new_slide(prs, "模型结构拆解：先压缩，再量化，再恢复", "Attention Projection 负责通道变换，MCQ 负责离散容量扩展。", "这一页把模型图拆成中文流程。输入图像经过 encoder 变成 B、N、C 的 token 序列；compression 把 C 压到 c，降低量化难度；MCQ 把每个 token 拆成多个 chunk，用多个码本生成 B、M、N 的离散 indices；expansion 再把通道恢复到 decoder 需要的维度。生成侧用 decoder 重建图像，理解侧用 pooled feature 和文本特征做对比学习。")
    add_flow(slide, ["图像", "编码器", "通道压缩", "MCQ 多码本", "通道恢复", "解码器", "重建图"], y=1.75, x0=0.65, box_w=1.52, gap=0.34)
    add_card(slide, "输入 token", "[B, 3, H, W]\n-> [B, N, C]", 0.9, 3.25, 2.55, 1.35, BLUE, body_size=15)
    add_card(slide, "量化前", "[B, N, C]\n-> [B, N, c]", 3.65, 3.25, 2.55, 1.35, GREEN, body_size=15)
    add_card(slide, "离散 token", "indices\n[B, M, N]", 6.4, 3.25, 2.55, 1.35, ORANGE, body_size=15)
    add_card(slide, "输出接口", "重建图像\n理解特征", 9.15, 3.25, 2.55, 1.35, BLUE, body_size=15)
    add_text(slide, "论文模块对应：视觉编码器 / 注意力投影 / 多码本量化 / 视觉解码器", 0.95, 5.25, 11.45, 0.35, size=17, color=DARK, align=PP_ALIGN.CENTER)
    add_conclusion(slide, "这套结构解决的是离散 token 容量不足，而不只是多加一个 loss")

    slide = new_slide(prs, "UniTok 总体框架", "图像经过通道压缩、MCQ、通道恢复后重建，并保留理解接口。", "这一页是 UniTok 的整体框架。图像进入 encoder 得到 token，通道压缩后进入 MCQ，MCQ 产生离散 indices，再经过通道恢复和 decoder 重建图像。理解侧可以使用视觉特征和文本特征做对比学习。")
    add_flow(slide, ["图像", "编码器", "通道压缩", "MCQ", "通道恢复", "解码器", "重建图"], y=2.25, x0=0.7, box_w=1.42, gap=0.34)
    add_card(slide, "生成路径", "indices -> 解码器 -> 图像", 1.55, 4.2, 3.5, 1.15, BLUE)
    add_card(slide, "理解路径", "视觉特征 -> 文本对齐", 5.3, 4.2, 3.5, 1.15, GREEN)
    add_card(slide, "统一 token", "[B, M, N] 离散索引", 9.05, 4.2, 3.5, 1.15, ORANGE)
    add_conclusion(slide, "同一套 token 同时连接生成与理解")

    slide = new_slide(prs, "训练目标：完整论文 vs tiny 复现", "本项目保留核心训练项，省略高成本感知/对抗/大规模对比训练。", "论文完整目标包含像素重建、VQ、感知、对抗和图文对比。我的 Jittor tiny 复现默认使用重建损失加 VQ 损失，因为这样能在普通环境中稳定跑通，并验证 tokenizer 的核心结构。")
    add_card(slide, "论文完整目标", "L_R + L_VQ + L_P + L_G + L_contra", 1.0, 1.6, 5.4, 1.45, BLUE)
    add_card(slide, "本项目默认目标", "L_R + λ_VQ L_VQ", 7.0, 1.6, 5.0, 1.45, GREEN)
    add_table(slide, [["项", "含义", "tiny 是否默认"], ["L_R", "像素重建", "是"], ["L_VQ", "码本训练", "是"], ["L_P/L_G", "感知/对抗", "否"], ["L_contra", "图文语义", "接口保留"]], 1.25, 3.45, 10.8, 2.25, font_size=14)
    add_conclusion(slide, "复现目标：结构正确、流程跑通、证据完整")

    slide = new_slide(prs, "复现范围：Jittor tiny setting", "复现核心 tokenizer，不复现 DataComp-1B 级训练。", "这页说明边界。原论文依赖大规模数据和模型，本项目使用 64x64 demo 图像、小 CNN、少量 batch 训练。这样不能验证大规模 VQA 或 GenEval 指标，但可以验证 forward、encode/decode、训练 loss、重建可视化和 PyTorch/Jittor 模块对齐。")
    add_card(slide, "保留", "MCQ\nAttention Projection\nencode/decode\n训练与可视化", 1.0, 1.6, 4.8, 2.3, BLUE)
    add_card(slide, "轻量替代", "ViTamin-L -> tiny CNN\nDataComp-1B -> demo data\nrFID -> MSE/PSNR", 7.1, 1.6, 4.8, 2.3, ORANGE)
    add_conclusion(slide, "tiny setting 验证算法链路，不声称复现论文大规模指标")

    slide = new_slide(prs, "官方 PyTorch 到 Jittor 的映射", "迁移重点是模块职责和 shape 对齐。", "我先分析官方 PyTorch 代码，再在 Jittor 项目中建立对应文件。unitok.py 对应 tokenizer.py，quant.py 对应 mcq.py，attention projection 对应 attention_projection.py，大 backbone 用 tiny encoder/decoder 替代。")
    add_table(slide, [["PyTorch 官方", "Jittor 复现", "对齐重点"], ["models/unitok.py", "models/tokenizer.py", "端到端 graph"], ["models/quant.py", "models/mcq.py", "MCQ split/lookup/concat"], ["models/vqvae.py", "attention_projection.py", "compression/expansion"], ["trainer/loss", "engine + losses.py", "tiny training"]], 0.75, 1.45, 11.9, 3.7, font_size=14)
    add_conclusion(slide, "Jittor 版本对齐结构与接口，而非复制大模型参数")

    slide = new_slide(prs, "Jittor 代码实现详解", "下面进入独立代码讲解，覆盖 tokenizer、MCQ、Attention、loss、训练和对齐。", "接下来进入独立的代码实现详解。这里不是论文和代码合计 10 分钟，而是单独给代码实现约 10 分钟。16 页代码页每页预计 40 秒，每页都有路径、论文模块、shape、关键代码和一句话总结。其中 losses.py 被拆成两页，完整展示重建损失、VQ 损失和对比损失的定义。")
    add_text(slide, "Jittor 代码实现详解", 1.0, 1.35, 11.3, 0.75, size=34, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "从 tokenizer 到训练评估的完整实现链路", 2.2, 2.35, 8.9, 0.45, size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_flow(slide, ["Tokenizer", "MCQ", "Attention", "Backbone", "Loss", "Train/Eval", "Align"], y=3.65, x0=0.8, box_w=1.45, gap=0.35)
    add_conclusion(slide, "主页面看代码骨架，详细解释在备注和讲稿")

    for spec in CODE_SLIDES:
        add_code_slide(prs, spec)

    slide = new_slide(prs, "项目如何运行：一键复现实验", "按 README 命令顺序运行，可生成训练、评估、曲线和 PPT 材料。", "这一页回答项目如何运行。推荐先进入项目目录并激活 unitok_jittor 环境。如果还没有数据，就运行 prepare_demo_data 生成 demo 图像。然后执行 train_tokenizer 训练 tiny tokenizer，eval_reconstruction 做重建评估，plot_loss 生成 loss 曲线，compare_with_pytorch 生成对齐报告，最后 build_ppt 重新生成答辩材料。所有输出都在 outputs 目录下。")
    add_tag(slide, "环境：conda activate unitok_jittor", 0.75, 1.05, 5.2, BLUE, size=12)
    add_tag(slide, "目录：D:\\newproject\\jittor-unitok", 6.15, 1.05, 5.9, GREEN, size=12)
    add_code(slide, """cd D:\\newproject\\jittor-unitok
conda activate unitok_jittor

python scripts\\prepare_demo_data.py --output data/demo --image-size 64
python -m jittor_unitok.engine.train_tokenizer --config configs/unitok_demo.yaml --epochs 1 --batch-size 2 --no-cuda
python -m jittor_unitok.engine.eval_reconstruction --checkpoint outputs/checkpoints/unitok_last.pkl --data-root data/demo --output-dir outputs --no-cuda
python scripts\\plot_loss.py --csv outputs/curves/loss.csv --output outputs/curves/loss_curve.png
python scripts\\generate_ppt_figures.py --ma-window 500
python scripts\\build_ppt.py""", 0.78, 1.65, 11.75, 3.95, size=10.5)
    add_conclusion(slide, "运行结果：outputs/logs、outputs/curves、outputs/reconstructions、outputs/ppt")

    slide = new_slide(prs, "结果：30 轮 loss 整体下降", "完整 loss.csv 覆盖 30 轮训练，最终 step 达到 187500。", "这一页展示 30 轮完整训练的 loss 证据。这里使用 outputs/ppt/figures/loss_epoch_30.png，它是从完整 loss.csv 按 epoch 聚合得到的，不是早期 demo 的 1 到 8 step 曲线。可以看到 total loss、重建损失和 VQ 损失整体下降，说明 Jittor 版 tokenizer 的训练流程持续有效。")
    curve = ROOT / "outputs" / "ppt" / "figures" / "loss_epoch_30.png"
    if not curve.exists():
        curve = ROOT / "outputs" / "ppt" / "figures" / "loss_curve_ma.png"
    if not curve.exists():
        curve = ROOT / "outputs" / "curves" / "loss_curve.png"
    picture(slide, curve, 0.62, 1.25, 6.55)
    add_card(slide, "训练规模", read_training_manifest(), 7.55, 1.25, 4.85, 0.95, BLUE, body_size=15)
    add_card(slide, "epoch 平均 total loss", read_epoch_loss_summary(), 7.55, 2.45, 4.85, 1.35, GREEN, body_size=14)
    add_card(slide, "train.log / loss.csv", read_final_train_summary() + "\n" + read_loss_tail(), 7.55, 4.12, 4.85, 1.45, BLUE, body_size=11)
    add_conclusion(slide, "30 轮训练的 epoch 平均 loss 明显下降，最后 step 达到 187500")

    slide = new_slide(prs, "结果：重建图验证 encode-decode", "重建图保留主要颜色和结构，说明离散 token 流程有效。", "这一页展示更清晰的重建对比图。左列是输入原图，中列是 tokenizer encode-decode 之后的重建图，右列是误差图。CIFAR-10 原始分辨率只有 32x32，所以细节不会很锐利；但主体颜色和轮廓能够恢复，说明 Jittor 版 encoder、MCQ 和 decoder 链路是有效的。")
    showcase = ROOT / "outputs" / "ppt" / "figures" / "reconstruction_showcase_clear.png"
    if not picture(slide, showcase, 0.55, 1.08, 9.25):
        picture(slide, ROOT / "outputs" / "reconstructions" / "train_epoch_30.png", 0.75, 1.35, 5.75)
        picture(slide, ROOT / "outputs" / "reconstructions" / "eval_batch_1.png", 6.85, 1.35, 5.75)
    metrics = read_metrics()
    add_card(slide, "整体评估", f"MSE {metrics.get('mse', 0):.4f}\nL1 {metrics.get('l1', 0):.4f}\nPSNR {metrics.get('psnr', 0):.2f} dB", 10.05, 1.45, 2.55, 1.65, GREEN, body_size=14)
    add_card(slide, "读图方式", "左：原图\n中：重建\n右：误差图", 10.05, 3.55, 2.55, 1.45, BLUE, body_size=14)
    add_conclusion(slide, "主体颜色和轮廓可重建，验证 encode-decode 流程有效", y=6.45)

    slide = new_slide(prs, "结果：PyTorch/Jittor 对齐", "模块、shape、训练产物三层对齐，支撑结构复现结论。", "对齐结果不是完整大模型数值复现，而是结构和工程行为对齐。模块上有对应文件，shape 上保持 image、latent、indices、reconstruction 的一致，行为上 quick_start、pytest、train 和 eval 都能跑通。")
    add_table(slide, [["对齐维度", "结论"], ["模块职责", "tokenizer / MCQ / attention / train 均有对应"], ["关键 shape", "image [B,3,H,W]；indices [B,M,N]；recon [B,3,H,W]"], ["运行行为", "quick_start / pytest / train / eval 已跑通"], ["训练趋势", "loss.csv + train.log 可复现"]], 0.85, 1.35, 11.6, 3.8, font_size=15)
    add_conclusion(slide, "Jittor 版本完成结构级与流程级对齐")

    slide = new_slide(prs, "结果：实验材料可复现", "从数据到 PPT 的产物链路完整。", "这一页强调作业可检查性。准备数据、训练、评估、画曲线、生成 PPT，都有命令和输出路径。结果不依赖手工截图，而是项目脚本生成。")
    add_flow(slide, ["prepare", "train", "eval", "plot", "compare", "PPT"], y=2.2, x0=1.25, box_w=1.55, gap=0.48)
    add_table(slide, [["产物", "路径"], ["train.log", "outputs/logs/train.log"], ["30 epoch loss", "outputs/ppt/figures/loss_epoch_30.png"], ["recon", "outputs/reconstructions/*.png"], ["slides", "outputs/ppt/*.pptx/pdf/md"]], 1.4, 3.55, 10.5, 2.1, font_size=14)
    add_conclusion(slide, "实验、可视化和答辩材料由同一项目生成")

    slide = new_slide(prs, "我的思考：统一 token 的核心矛盾", "统一生成与理解的关键，是共享 token 空间中的容量和语义表达平衡。", "我认为这个方向的核心矛盾不是简单多加一个 CLIP loss，而是共享 token 空间的容量问题。生成需要细节，理解需要语义，离散 token 要同时容纳这两类信息。UniTok 的贡献在于从 tokenizer 容量角度处理瓶颈，而不是只堆训练目标。")
    add_card(slide, "领域矛盾", "共享 token 空间\n容量有限 + 语义复杂", 1.0, 1.65, 3.6, 1.8, ORANGE)
    add_card(slide, "UniTok 贡献", "MCQ 扩组合容量\nAttention 保上下文", 4.9, 1.65, 3.6, 1.8, BLUE)
    add_card(slide, "我的判断", "结构创新比简单堆 loss 更关键", 8.8, 1.65, 3.6, 1.8, GREEN)
    add_conclusion(slide, "统一 tokenizer 的上限取决于离散表达能力")

    slide = new_slide(prs, "我的思考：局限、下一步与应用", "tiny 复现验证流程，真实价值还需要理解任务和应用场景检验。", "这个项目的局限也很明确：完整训练成本极高，tiny 复现不能证明大规模 VQA、GenEval 或图文检索指标。下一步我想尝试更轻量的 MCQ、蒸馏 CLIP 语义，并在小型 VQA 或图文检索数据上验证理解能力。应用上可以考虑多轮图像编辑、工业检测报告生成、医疗影像辅助分析和统一多模态助手。")
    add_table(slide, [["方向", "关键词"], ["局限", "训练成本高；tiny 不能验证大规模指标"], ["下一步", "轻量 MCQ；CLIP 蒸馏；小型 VQA/检索"], ["应用", "图像编辑；工业报告；医疗影像；多模态助手"]], 1.0, 1.55, 11.3, 3.15, font_size=16)
    add_conclusion(slide, "从 tiny 复现走向真实应用，需要补理解侧验证")

    slide = new_slide(prs, "资源受限说明与 GitHub 材料", "本项目是核心结构复现，不是论文级大规模训练。", "原论文使用 DataComp-1B、大规模视觉和语言模型、长时间训练。本项目使用小数据、小模型和少量 batch 验证核心思想。README 已记录环境、数据、训练、评估、可视化、PPT 和 GitHub 占位。")
    add_card(slide, "原论文", "DataComp-1B\n大模型\n长训练", 1.0, 1.55, 3.3, 1.7, BLUE)
    add_card(slide, "本项目", "CIFAR-10\n小型 CNN\n核心流程", 5.0, 1.55, 3.3, 1.7, GREEN)
    add_card(slide, "提交材料", "README\nPPT/PDF\nSpeech.md", 9.0, 1.55, 3.3, 1.7, ORANGE)
    add_text(slide, "GitHub 占位： https://github.com/<your-name>/jittor-unitok", 1.25, 4.45, 10.9, 0.4, size=20, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, "最终命名预留：姓名-播种期.pdf", 3.1, 5.1, 7.0, 0.35, size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_conclusion(slide, "边界讲清楚，复现证据链完整")

    slide = new_slide(prs, "总结", "Jittor 版 UniTok tokenizer 已完成核心结构复现与展示材料生成。", "总结一下，本项目完成了 Jittor 版 UniTok tokenizer 的核心结构，包括 MCQ、Attention Projection、forward、encode、decode、训练、评估和可视化。实验结果证明 tiny setting 下流程跑通。后续可以继续向真实图文数据、CLIP 蒸馏和轻量 MCQ 方向扩展。")
    add_card(slide, "实现", "MCQ\nAttention Projection\nTokenizer API", 1.0, 1.6, 3.5, 1.9, BLUE)
    add_card(slide, "验证", "loss 曲线\n重建图\nPyTorch/Jittor 对齐", 4.9, 1.6, 3.5, 1.9, GREEN)
    add_card(slide, "后续", "真实数据\n理解任务\n轻量化 MCQ", 8.8, 1.6, 3.5, 1.9, ORANGE)
    add_conclusion(slide, "完成的是可运行、可展示、边界清晰的核心复现")

    return prs


def write_speech() -> None:
    parts = [
        "# UniTok Jittor Reproduction 逐页演讲稿",
        "",
        "说明：主页面只保留关键词、图表和结论句；详细讲解写在本文件和 PPT 演讲者备注中。",
        "代码讲解章节包含 1 页过渡页和 16 个代码内容页；代码内容页按 40 秒/页设计，总计约 640 秒，即约 10 分钟多一点。",
        "",
    ]
    for note in NOTES:
        parts.append(f"## 第 {note['page']} 页：{note['title']}")
        parts.append("")
        parts.append(f"**本页核心观点：** {note['core']}")
        parts.append("")
        parts.append(f"**讲稿正文：** {note['body']}")
        parts.append("")
    SPEECH_OUT.write_text("\n".join(parts), encoding="utf-8")
    NOTES_JSON_OUT.write_text(json.dumps(NOTES, ensure_ascii=False, indent=2), encoding="utf-8")


def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(PPTX_OUT)
    write_speech()
    print(f"saved {PPTX_OUT}")
    print(f"saved {SPEECH_OUT}")
    print(f"saved {NOTES_JSON_OUT}")


if __name__ == "__main__":
    main()
